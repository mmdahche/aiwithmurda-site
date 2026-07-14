"""Tests for the to-markdown skill (doc/data -> Markdown, fail-closed, no network).

Drives the REAL markitdown via the production API ``convert_path`` /
``convert_bytes`` on tiny in-test fixtures (HTML, CSV, XLSX, PPTX, DOCX, and
-- when reportlab is on the machine -- a real PDF). Pure structural-preservation
assertions.

No-false-green guards:
  * fail-closed paths return non-zero and write NOTHING (oversized, corrupt,
    disallowed extension, URL input);
  * the scrubber actually runs on output (recording fake scrubber asserts
    the call happened on every public entry point);
  * a real-CLI subprocess smoke test exercises the production code path
    end-to-end.

The optional guardrails firewall is stubbed on-disk (in ``$GUARDRAILS_HOME``)
for the tests that assert the "default scrubber IS the guard" wiring works.

Hermetic: all fixtures are built inside the test's tmp_path. If markitdown is
somehow not installed, every test is skipped with a clear message rather than
passing silently.
"""
from __future__ import annotations

import importlib.util
import io
import os
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest

# Load the module under test (file-path import).
TM_DIR = Path(__file__).resolve().parents[1]
TM_PATH = TM_DIR / "tomarkdown.py"


def _load(name: str = "tomarkdown_under_test"):
    spec = importlib.util.spec_from_file_location(name, str(TM_PATH))
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


tm = _load()

# Skip the whole suite if markitdown is not importable.
markitdown_available = importlib.util.find_spec("markitdown") is not None
reportlab_available = importlib.util.find_spec("reportlab") is not None
pytestmark = pytest.mark.skipif(
    not markitdown_available,
    reason="markitdown not installed (pip install 'markitdown[pdf,docx,pptx,xlsx]')")

# A fake token the stub firewall will recognise -- built by concat so this file
# contains no contiguous secret-shaped literal.
TOKEN_LOOKALIKE = "abc_" + "1234567890ABCDEFGHIJ"


@pytest.fixture(autouse=True)
def _isolate_guardrails(monkeypatch):
    """Reset the module-level guardrails cache and clear any GUARDRAILS_HOME."""
    monkeypatch.delenv("GUARDRAILS_HOME", raising=False)
    monkeypatch.setattr(tm, "_FW", None)
    monkeypatch.setattr(tm, "_FW_ERR", None)
    monkeypatch.setattr(tm, "_FW_WARNED", False)


def _install_stub_firewall(monkeypatch, tmp_path, marker="[REDACTED]"):
    """Drop a tiny stub firewall inside a fake ``$GUARDRAILS_HOME`` and point
    the module at it. The stub redacts our TOKEN_LOOKALIKE pattern."""
    root = tmp_path / "guardrails"
    fw = root / "redaction-firewall" / "firewall.py"
    fw.parent.mkdir(parents=True, exist_ok=True)
    fw.write_text(
        "def guard(text):\n"
        "    if text is None:\n"
        "        return ''\n"
        "    import re\n"
        "    return re.sub(r'[A-Za-z]{3,}_[A-Za-z0-9]{10,}',\n"
        "                  '" + marker + "', text)\n",
        encoding="utf-8",
    )
    monkeypatch.setenv("GUARDRAILS_HOME", str(root))
    monkeypatch.setattr(tm, "_FW", None)
    monkeypatch.setattr(tm, "_FW_ERR", None)
    monkeypatch.setattr(tm, "_FW_WARNED", False)
    return fw


# ── fixture builders ────────────────────────────────────────────────────────────

def _write_html(path: Path) -> Path:
    path.write_bytes(
        b"<html><head><title>T</title></head>"
        b"<body><h1>Heading One</h1>"
        b"<p>A paragraph of text.</p>"
        b"<ul><li>alpha</li><li>beta</li></ul></body></html>")
    return path


def _write_csv(path: Path) -> Path:
    path.write_bytes(b"name,age,role\nAlice,30,engineer\nBob,25,designer\n")
    return path


def _write_xlsx(path: Path) -> Path:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["name", "age", "role"])
    ws.append(["Alice", 30, "engineer"])
    ws.append(["Bob", 25, "designer"])
    wb.save(str(path))
    return path


def _write_pptx(path: Path) -> Path:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide Heading"
    slide.placeholders[1].text = "Slide subtitle line"
    prs.save(str(path))
    return path


def _build_minimal_docx_bytes() -> bytes:
    """A hand-built minimal Office Open XML wordprocessing doc (a 3-file zip)."""
    content_types = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        b'<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        b'<Default Extension="xml" ContentType="application/xml"/>'
        b'<Override PartName="/word/document.xml" '
        b'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        b'</Types>')
    rels = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        b'<Relationship Id="rId1" '
        b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        b'Target="word/document.xml"/></Relationships>')
    document = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        b'<w:body>'
        b'<w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr>'
        b'<w:r><w:t>Hello Doc</w:t></w:r></w:p>'
        b'<w:p><w:r><w:t>This is a paragraph of body text.</w:t></w:r></w:p>'
        b'</w:body></w:document>')
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document)
    return buf.getvalue()


def _write_docx(path: Path) -> Path:
    path.write_bytes(_build_minimal_docx_bytes())
    return path


def _write_pdf(path: Path) -> Path:
    from reportlab.pdfgen.canvas import Canvas  # type: ignore
    c = Canvas(str(path))
    c.setFont("Helvetica", 14)
    c.drawString(72, 720, "Hello PDF World")
    c.drawString(72, 700, "second line of text")
    c.showPage()
    c.save()
    return path


class RecordingScrubber:
    """Captures every text it was asked to scrub. Returns ``text + sentinel``."""

    SENTINEL = "\n<!-- scrubbed-by-test -->"

    def __init__(self):
        self.calls = []

    def __call__(self, text):
        self.calls.append(text)
        return (text or "") + self.SENTINEL


# ── 1) HAPPY-PATH conversions (structure preserved, scrubber called) ────────────

def test_convert_path_html_preserves_structure(tmp_path):
    p = _write_html(tmp_path / "page.html")
    sc = RecordingScrubber()
    out = tm.convert_path(p, scrubber=sc)
    assert "Heading One" in out
    assert "# Heading One" in out
    assert "alpha" in out and "beta" in out
    assert sc.SENTINEL in out
    assert len(sc.calls) == 1


def test_convert_path_csv_becomes_markdown_table(tmp_path):
    p = _write_csv(tmp_path / "people.csv")
    sc = RecordingScrubber()
    out = tm.convert_path(p, scrubber=sc)
    assert "| name | age | role |" in out
    assert "| Alice | 30 | engineer |" in out
    assert sc.SENTINEL in out


def test_convert_path_xlsx_becomes_markdown_table(tmp_path):
    p = _write_xlsx(tmp_path / "people.xlsx")
    sc = RecordingScrubber()
    out = tm.convert_path(p, scrubber=sc)
    assert "name" in out and "Alice" in out
    assert "|" in out
    assert sc.SENTINEL in out


def test_convert_path_pptx_extracts_slide_text(tmp_path):
    p = _write_pptx(tmp_path / "deck.pptx")
    sc = RecordingScrubber()
    out = tm.convert_path(p, scrubber=sc)
    assert "Slide Heading" in out
    assert "Slide subtitle line" in out
    assert sc.SENTINEL in out


def test_convert_path_docx_extracts_heading_and_body(tmp_path):
    p = _write_docx(tmp_path / "note.docx")
    sc = RecordingScrubber()
    out = tm.convert_path(p, scrubber=sc)
    assert "Hello Doc" in out
    assert "This is a paragraph of body text." in out
    assert sc.SENTINEL in out


@pytest.mark.skipif(not reportlab_available,
                    reason="reportlab not installed; cannot build a real PDF fixture")
def test_convert_path_pdf_extracts_text(tmp_path):
    p = _write_pdf(tmp_path / "doc.pdf")
    sc = RecordingScrubber()
    out = tm.convert_path(p, scrubber=sc)
    assert "Hello PDF World" in out
    assert "second line of text" in out
    assert sc.SENTINEL in out


def test_convert_bytes_html_via_mimetype():
    sc = RecordingScrubber()
    out = tm.convert_bytes(b"<h1>Inline</h1><p>x</p>", "text/html", scrubber=sc)
    assert "# Inline" in out and sc.SENTINEL in out


# ── 2) FAIL-CLOSED validation (no lie-green) ────────────────────────────────────

def test_url_input_is_refused_https(tmp_path):
    with pytest.raises(tm.ToMarkdownError, match="URL inputs are REFUSED"):
        tm.convert_path("https://example.com/doc.html",
                        scrubber=RecordingScrubber())


def test_url_input_is_refused_http_case_insensitive():
    with pytest.raises(tm.ToMarkdownError, match="URL inputs are REFUSED"):
        tm.convert_path("HTTP://Example.COM/x.html",
                        scrubber=RecordingScrubber())


def test_url_input_is_refused_file_and_data():
    for url in ("file:///etc/passwd", "data:text/html,<h1>x</h1>", "ftp://h/x.txt"):
        with pytest.raises(tm.ToMarkdownError, match="URL inputs are REFUSED"):
            tm.convert_path(url, scrubber=RecordingScrubber())


def test_disallowed_extension_is_refused(tmp_path):
    p = tmp_path / "thing.exe"
    p.write_bytes(b"MZ\x90\x00")
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="extension '.exe' is NOT on the allowlist"):
        tm.convert_path(p, scrubber=sc)
    assert sc.calls == []


def test_missing_file_is_refused(tmp_path):
    with pytest.raises(tm.ToMarkdownError, match="does not exist"):
        tm.convert_path(tmp_path / "absent.html", scrubber=RecordingScrubber())


def test_empty_file_is_refused(tmp_path):
    p = tmp_path / "empty.html"
    p.write_bytes(b"")
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="is empty"):
        tm.convert_path(p, scrubber=sc)
    assert sc.calls == []


def test_oversized_file_is_refused(tmp_path):
    p = tmp_path / "big.html"
    p.write_bytes(b"<p>x</p>" * 10)
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="exceeds the .* cap"):
        tm.convert_path(p, scrubber=sc, max_bytes=10)
    assert sc.calls == []


def test_corrupt_zip_based_file_is_refused(tmp_path):
    p = tmp_path / "broken.docx"
    p.write_bytes(b"this is definitely not a zip")
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="not a valid zip-based file"):
        tm.convert_path(p, scrubber=sc)
    assert sc.calls == []


def test_zip_traversal_entry_is_refused(tmp_path):
    p = tmp_path / "evil.docx"
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("word/document.xml", b"<x/>")
        zf.writestr("../../etc/passwd", b"root:x:0:0:")
    p.write_bytes(buf.getvalue())
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="traversal"):
        tm.convert_path(p, scrubber=sc)
    assert sc.calls == []


def test_zip_bomb_uncompressed_cap_is_enforced(tmp_path):
    p = tmp_path / "bomb.epub"
    payload = b"\x00" * (tm.MAX_ZIP_UNCOMPRESSED + 1)
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("a.txt", payload)
    p.write_bytes(buf.getvalue())
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="zip-bomb guard"):
        tm.convert_path(p, scrubber=sc, max_bytes=10 * 1024 * 1024)
    assert sc.calls == []


def test_zip_extension_is_refused_in_v1(tmp_path):
    p = tmp_path / "anything.zip"
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("inner.txt", b"hi")
    p.write_bytes(buf.getvalue())
    sc = RecordingScrubber()
    with pytest.raises(tm.ToMarkdownError, match="extension '.zip' is NOT on the allowlist"):
        tm.convert_path(p, scrubber=sc)
    assert sc.calls == []


# ── 3) The OUTPUT SCRUBBER runs on every public entrypoint ──────────────────────

def test_scrubber_failure_propagates_and_no_text_returned(tmp_path):
    """If the scrubber raises, the caller MUST see the exception -- never the
    un-scrubbed Markdown."""
    p = _write_html(tmp_path / "page.html")

    class _Boom(Exception):
        pass

    def _scrub_explodes(_text):
        raise _Boom("redactor exploded")

    with pytest.raises(_Boom):
        tm.convert_path(p, scrubber=_scrub_explodes)


def test_default_scrubber_passthrough_when_firewall_missing(monkeypatch, tmp_path, capsys):
    """Without a firewall installed, the default scrubber passes text through
    unchanged and emits a one-time stderr note."""
    monkeypatch.delenv("GUARDRAILS_HOME", raising=False)
    empty_home = tmp_path / "empty-home"
    empty_home.mkdir()
    monkeypatch.setenv("HOME", str(empty_home))
    p = tmp_path / "leak.txt"
    p.write_bytes(("token is " + TOKEN_LOOKALIKE + "\nother text\n").encode())
    out = tm.convert_path(p)  # default scrubber
    assert TOKEN_LOOKALIKE in out
    err = capsys.readouterr().err
    assert "optional Safe-Autonomy Guardrails firewall not found" in err


def test_default_scrubber_uses_guard_when_firewall_installed(monkeypatch, tmp_path):
    """With a stub firewall installed, the default scrubber is the guardrails
    guard() and it actually redacts."""
    _install_stub_firewall(monkeypatch, tmp_path, marker="[REDACTED]")
    p = tmp_path / "leak.txt"
    p.write_bytes(("token is " + TOKEN_LOOKALIKE + "\nother text\n").encode())
    out = tm.convert_path(p)
    assert TOKEN_LOOKALIKE not in out
    assert "[REDACTED]" in out


# ── 4) REAL CLI smoke test (end-to-end, exercises the production code path) ─────

def _run(args, cwd):
    env = dict(os.environ)
    env["PYTHONUTF8"] = "1"
    env["PYTHONIOENCODING"] = "utf-8"
    return subprocess.run(
        [sys.executable, str(TM_PATH)] + args,
        cwd=str(cwd), capture_output=True, text=True, env=env)


def test_cli_convert_to_stdout(tmp_path):
    p = _write_html(tmp_path / "page.html")
    r = _run(["convert", "--in", str(p)], tmp_path)
    assert r.returncode == 0
    assert "Heading One" in r.stdout
    assert "# Heading One" in r.stdout


def test_cli_convert_to_outfile(tmp_path):
    p = _write_csv(tmp_path / "data.csv")
    out = tmp_path / "out.md"
    r = _run(["convert", "--in", str(p), "--out", str(out)], tmp_path)
    assert r.returncode == 0 and out.is_file()
    body = out.read_text(encoding="utf-8")
    assert "| Alice | 30 |" in body


def test_cli_refuses_url_input(tmp_path):
    r = _run(["convert", "--in", "https://example.com/x.html"], tmp_path)
    assert r.returncode == 2
    assert "URL inputs are REFUSED" in r.stderr


def test_cli_refuses_disallowed_extension_and_writes_nothing(tmp_path):
    p = tmp_path / "thing.exe"
    p.write_bytes(b"MZ\x90\x00")
    out = tmp_path / "should-not-exist.md"
    r = _run(["convert", "--in", str(p), "--out", str(out)], tmp_path)
    assert r.returncode == 2
    assert not out.exists()


def test_cli_refuses_oversized(tmp_path):
    p = _write_html(tmp_path / "page.html")
    r = _run(["convert", "--in", str(p), "--max-bytes", "5"], tmp_path)
    assert r.returncode == 2
    assert "exceeds" in r.stderr or "cap" in r.stderr


def test_cli_requires_subcommand(tmp_path):
    r = _run([], tmp_path)
    assert r.returncode != 0


# ── 5) REGRESSION TESTS ─────────────────────────────────────────────────────────

def test_regression_zip_bomb_guard_inspects_in_memory_bytes_not_path(tmp_path):
    """convert_path used to call ``_zip_bomb_guard(p, nbytes)`` which re-opened
    the file from disk, creating a verified-bytes != processed-bytes window.
    The fix passes the same bytes that ``_run_markitdown`` will consume."""
    benign = io.BytesIO()
    with zipfile.ZipFile(benign, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("word/document.xml", b"<x/>")
    p = tmp_path / "swap.docx"
    p.write_bytes(benign.getvalue())

    bomb = io.BytesIO()
    payload = b"\x00" * (tm.MAX_ZIP_UNCOMPRESSED + 1)
    with zipfile.ZipFile(bomb, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("a.txt", payload)
    bomb_bytes = bomb.getvalue()

    real_read_bytes = Path.read_bytes

    def lying_read_bytes(self):
        if self.name == "swap.docx":
            return bomb_bytes
        return real_read_bytes(self)

    original = Path.read_bytes
    Path.read_bytes = lying_read_bytes
    try:
        sc = RecordingScrubber()
        with pytest.raises(tm.ToMarkdownError, match="zip-bomb guard"):
            tm.convert_path(p, scrubber=sc)
        assert sc.calls == []
    finally:
        Path.read_bytes = original


def test_regression_size_cap_revalidated_after_read_bytes(tmp_path):
    """The size cap was checked against ``p.stat().st_size`` BEFORE
    ``read_bytes()``. If the file grew between those two calls, an oversized
    payload could reach the parser uncapped."""
    p = tmp_path / "page.html"
    p.write_bytes(b"<p>x</p>")

    real_read_bytes = Path.read_bytes
    huge = b"<p>" + (b"x" * 4096) + b"</p>"

    def lying_read_bytes(self):
        if self.name == "page.html":
            return huge
        return real_read_bytes(self)

    original = Path.read_bytes
    Path.read_bytes = lying_read_bytes
    try:
        sc = RecordingScrubber()
        with pytest.raises(tm.ToMarkdownError, match="post-read"):
            tm.convert_path(p, scrubber=sc, max_bytes=100)
        assert sc.calls == []
    finally:
        Path.read_bytes = original


def test_regression_convert_bytes_xml_mimetype_is_deterministic():
    """``convert_bytes`` used ``mimetypes.guess_extension`` which consults the
    OS / Windows registry and could return ``.xsl`` for ``application/xml``,
    falsely refusing a valid XML doc. The fix is a hardcoded MIME -> extension
    table that maps ONLY to allowlist entries."""
    sc = RecordingScrubber()
    xml = b'<?xml version="1.0"?><root><item>hello-xml</item></root>'
    out = tm.convert_bytes(xml, "application/xml", scrubber=sc)
    assert sc.SENTINEL in out
    assert len(sc.calls) == 1

    sc2 = RecordingScrubber()
    out2 = tm.convert_bytes(b"line one\nline two\n",
                            "text/plain; charset=utf-8", scrubber=sc2)
    assert sc2.SENTINEL in out2
    with pytest.raises(tm.ToMarkdownError, match="NOT on the allowlist"):
        tm.convert_bytes(b"x", "application/x-totally-made-up",
                         scrubber=RecordingScrubber())
